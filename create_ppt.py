from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# 프레젠테이션 생성
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# 색상 정의 (딥네이비, 화이트, 블랙)
DEEP_NAVY = RGBColor(20, 40, 80)  # #142850
NAVY_LIGHT = RGBColor(39, 60, 117)  # #273C75
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(128, 128, 128)

def add_blank_slide():
    """빈 슬라이드 추가"""
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_title_box(slide, text, top, left, width, height, font_size, bold=True, color=WHITE):
    """제목 텍스트박스 추가"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    return textbox

def add_text_box(slide, text, top, left, width, height, font_size, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    """일반 텍스트박스 추가"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.line_spacing = 1.3
    return textbox

def add_bullet_text(slide, items, top, left, width, height, font_size=16, color=WHITE):
    """불릿 포인트 추가"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(items):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.level = 0
        p.line_spacing = 1.4

    return textbox

# 슬라이드 1: 표지
slide1 = add_blank_slide()
slide1.background.fill.solid()
slide1.background.fill.fore_color.rgb = DEEP_NAVY

add_title_box(slide1, "광고·릴스·콘텐츠의", 2.0, 1, 8, 0.8, 36, color=WHITE)
add_title_box(slide1, "국룰 구조 PRD", 2.7, 1, 8, 1, 54, color=WHITE)
add_text_box(slide1, "타겟 → 오퍼 → 크리에이티브 → 랜딩 → 전환", 4.2, 1, 8, 0.5, 18, color=GRAY, align=PP_ALIGN.CENTER)

# 슬라이드 2: 목적
slide2 = add_blank_slide()
slide2.background.fill.solid()
slide2.background.fill.fore_color.rgb = WHITE

# 제목 배경
title_bg = slide2.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = DEEP_NAVY
title_bg.line.fill.background()

add_title_box(slide2, "이 문서의 목적", 0.3, 1, 8, 0.6, 36, color=WHITE)

add_text_box(slide2, "이 문서는 다음과 같은 분들을 위한 가이드입니다:", 1.8, 1, 8, 0.5, 20, bold=True, color=BLACK, align=PP_ALIGN.LEFT)

add_bullet_text(slide2, [
    "광고를 처음 접하는 원장님",
    "릴스는 만들지만 왜 안 되는지 모르는 원장님",
    "마케팅을 '감'으로만 해왔던 원장님"
], 2.6, 1.5, 7, 1.5, 18, color=BLACK)

# 하단 강조 박스
emphasis_box = slide2.shapes.add_shape(1, Inches(1), Inches(5.0), Inches(8), Inches(1.8))
emphasis_box.fill.solid()
emphasis_box.fill.fore_color.rgb = NAVY_LIGHT
emphasis_box.line.fill.background()

add_text_box(slide2, "이 구조는 유료 광고뿐만 아니라", 5.2, 1.5, 7, 0.4, 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide2, "릴스, 블로그, 상세페이지, 상담 구조 등", 5.6, 1.5, 7, 0.4, 18, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide2, "모든 고객 유입 콘텐츠에 공통으로 적용됩니다", 6.0, 1.5, 7, 0.4, 18, color=WHITE, align=PP_ALIGN.CENTER)

# 슬라이드 3: 타겟 (1)
slide3 = add_blank_slide()
slide3.background.fill.solid()
slide3.background.fill.fore_color.rgb = WHITE

title_bg = slide3.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = DEEP_NAVY
title_bg.line.fill.background()

add_title_box(slide3, "1️⃣ 타겟 (Target)", 0.3, 1, 8, 0.6, 36, color=WHITE)

# 정의 박스
def_box = slide3.shapes.add_shape(1, Inches(1), Inches(1.6), Inches(8), Inches(1.2))
def_box.fill.solid()
def_box.fill.fore_color.rgb = NAVY_LIGHT
def_box.line.fill.background()

add_text_box(slide3, '"내가 팔고 싶은 사람"이 아니라', 1.8, 1.5, 7, 0.4, 20, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide3, '"실제로 돈을 쓰고, 선택을 하는 사람"', 2.2, 1.5, 7, 0.4, 20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide3, "❌ 원장님들이 흔히 하는 착각", 3.2, 1, 8, 0.4, 18, bold=True, color=BLACK, align=PP_ALIGN.LEFT)

add_bullet_text(slide3, [
    '"20~30대 받고 싶어요"',
    '"젊은 손님이 많아야 잘 되는 거 아닌가요?"',
    '"요즘은 인스타니까 젊은 층이죠"'
], 3.7, 1.5, 7, 1.2, 16, color=GRAY)

add_text_box(slide3, "→ 이건 희망사항이지 타겟이 아닙니다", 5.0, 1.5, 7, 0.4, 18, bold=True, color=BLACK, align=PP_ALIGN.LEFT)

# 핵심 원칙
principle_box = slide3.shapes.add_shape(1, Inches(1), Inches(5.8), Inches(8), Inches(1.2))
principle_box.fill.solid()
principle_box.fill.fore_color.rgb = BLACK
principle_box.line.fill.background()

add_text_box(slide3, "💡 핵심 원칙", 6.0, 1.5, 7, 0.4, 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide3, "타겟이 정해지지 않으면 광고/릴스/콘텐츠는 전부 도박이 됩니다", 6.4, 1.5, 7, 0.4, 16, color=WHITE, align=PP_ALIGN.CENTER)

# 슬라이드 4: 타겟 (2) - 올바른 타겟 설정
slide4 = add_blank_slide()
slide4.background.fill.solid()
slide4.background.fill.fore_color.rgb = WHITE

title_bg = slide4.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = DEEP_NAVY
title_bg.line.fill.background()

add_title_box(slide4, "올바른 타겟 설정의 출발점", 0.3, 1, 8, 0.6, 32, color=WHITE)

add_text_box(slide4, "타겟은 아래 질문에서 나옵니다", 1.8, 1.5, 7, 0.4, 20, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

# 질문 박스들
questions = [
    "내 샵에서 실제로 돈을 잘 쓰는 사람은 누구인가?",
    "그 사람들은 왜 나를 선택했는가?",
    "그 사람들은 어디에서 나를 알게 되었는가?",
    "그 사람들은 가격에 어떤 반응을 보였는가?"
]

y_position = 2.6
for i, question in enumerate(questions):
    q_box = slide4.shapes.add_shape(1, Inches(1.5), Inches(y_position), Inches(7), Inches(0.7))
    q_box.fill.solid()
    q_box.fill.fore_color.rgb = NAVY_LIGHT if i % 2 == 0 else RGBColor(59, 80, 137)
    q_box.line.fill.background()

    add_text_box(slide4, question, y_position + 0.15, 1.8, 6.4, 0.4, 16, color=WHITE, align=PP_ALIGN.LEFT)
    y_position += 0.9

# 슬라이드 5: 오퍼 (1)
slide5 = add_blank_slide()
slide5.background.fill.solid()
slide5.background.fill.fore_color.rgb = WHITE

title_bg = slide5.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = DEEP_NAVY
title_bg.line.fill.background()

add_title_box(slide5, "2️⃣ 오퍼 (Offer)", 0.3, 1, 8, 0.6, 36, color=WHITE)

# 정의 박스
def_box = slide5.shapes.add_shape(1, Inches(1), Inches(1.6), Inches(8), Inches(1.2))
def_box.fill.solid()
def_box.fill.fore_color.rgb = NAVY_LIGHT
def_box.line.fill.background()

add_text_box(slide5, '"왜 지금, 왜 이걸, 왜 여기서', 1.8, 1.5, 7, 0.4, 20, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide5, '해야 하는지에 대한 이유"', 2.2, 1.5, 7, 0.4, 20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide5, "❌ 원장님들이 오해하는 오퍼", 3.2, 1, 8, 0.4, 18, bold=True, color=BLACK, align=PP_ALIGN.LEFT)

mistake_box = slide5.shapes.add_shape(1, Inches(1.5), Inches(3.7), Inches(7), Inches(0.9))
mistake_box.fill.solid()
mistake_box.fill.fore_color.rgb = RGBColor(240, 240, 240)
mistake_box.line.fill.background()

add_text_box(slide5, "할인 / 이벤트 / 리터치 포함", 3.85, 2, 6, 0.4, 18, color=GRAY, align=PP_ALIGN.CENTER)
add_text_box(slide5, "→ 이건 조건이지 오퍼가 아닙니다", 4.2, 2, 6, 0.4, 16, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

# 진짜 오퍼 박스
real_box = slide5.shapes.add_shape(1, Inches(1), Inches(5.0), Inches(8), Inches(1.6))
real_box.fill.solid()
real_box.fill.fore_color.rgb = BLACK
real_box.line.fill.background()

add_text_box(slide5, "✅ 진짜 오퍼의 역할", 5.2, 1.5, 7, 0.4, 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide5, '"관심은 있는데… 지금 당장 해야 할 이유는 모르겠어"', 5.6, 1.5, 7, 0.4, 16, color=GRAY, align=PP_ALIGN.CENTER)
add_text_box(slide5, "→ 오퍼는 이 망설임을 끊어주는 장치입니다", 6.0, 1.5, 7, 0.4, 18, color=WHITE, align=PP_ALIGN.CENTER)

# 슬라이드 6: 오퍼 (2) - 국룰 구성
slide6 = add_blank_slide()
slide6.background.fill.solid()
slide6.background.fill.fore_color.rgb = WHITE

title_bg = slide6.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = DEEP_NAVY
title_bg.line.fill.background()

add_title_box(slide6, "오퍼의 국룰 구성 (4요소)", 0.3, 1, 8, 0.6, 32, color=WHITE)

# 4요소 박스들
elements = [
    ("1. 대상 명확화", [
        '"눈썹 처음 하시는 분"',
        '"기존 시술 실패 경험 있으신 분"',
        '"자연눈썹 원하시는 40대 이상 고객"'
    ]),
    ("2. 지금 해야 하는 이유", [
        "기간 한정 / 선착순",
        "현재 상태(지금 안 하면 더 어려워짐)"
    ])
]

y_pos = 1.6
for title, items in elements:
    element_box = slide6.shapes.add_shape(1, Inches(0.8), Inches(y_pos), Inches(4.2), Inches(1.8))
    element_box.fill.solid()
    element_box.fill.fore_color.rgb = NAVY_LIGHT
    element_box.line.fill.background()

    add_text_box(slide6, title, y_pos + 0.1, 1, 4, 0.4, 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    item_text = "\n".join(items)
    add_text_box(slide6, item_text, y_pos + 0.6, 1.2, 3.6, 1, 14, color=WHITE, align=PP_ALIGN.LEFT)

    y_pos += 0

elements2 = [
    ("3. 리스크 제거", [
        "상담만 가능 / 강요 없음",
        "내 상태 진단 위주"
    ]),
    ("4. 얻는 결과 한 문장", [
        '"실패 확률을 줄이는 선택"',
        '"내 눈썹에 맞는 방향을 알게 됨"'
    ])
]

y_pos = 1.6
for title, items in elements2:
    element_box = slide6.shapes.add_shape(1, Inches(5.0), Inches(y_pos), Inches(4.2), Inches(1.8))
    element_box.fill.solid()
    element_box.fill.fore_color.rgb = RGBColor(59, 80, 137)
    element_box.line.fill.background()

    add_text_box(slide6, title, y_pos + 0.1, 5.2, 4, 0.4, 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    item_text = "\n".join(items)
    add_text_box(slide6, item_text, y_pos + 0.6, 5.4, 3.6, 1, 14, color=WHITE, align=PP_ALIGN.LEFT)

    y_pos += 0

# 하단 핵심 메시지
bottom_box = slide6.shapes.add_shape(1, Inches(1), Inches(5.8), Inches(8), Inches(1.2))
bottom_box.fill.solid()
bottom_box.fill.fore_color.rgb = BLACK
bottom_box.line.fill.background()

add_text_box(slide6, "💡 오퍼가 없으면 사람은 저장만 하고 행동하지 않습니다", 6.1, 1.5, 7, 0.4, 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 슬라이드 7: 크리에이티브 (1)
slide7 = add_blank_slide()
slide7.background.fill.solid()
slide7.background.fill.fore_color.rgb = WHITE

title_bg = slide7.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = DEEP_NAVY
title_bg.line.fill.background()

add_title_box(slide7, "3️⃣ 크리에이티브 (Creative)", 0.3, 1, 8, 0.6, 32, color=WHITE)

# 정의 박스
def_box = slide7.shapes.add_shape(1, Inches(1), Inches(1.6), Inches(8), Inches(1.2))
def_box.fill.solid()
def_box.fill.fore_color.rgb = NAVY_LIGHT
def_box.line.fill.background()

add_text_box(slide7, '광고/릴스에서 고객이 처음 마주치는 "입구"', 1.8, 1.5, 7, 0.4, 20, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide7, '목적: "이건 내 얘기다"라고 느끼게 하는 것', 2.2, 1.5, 7, 0.4, 20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide7, "❌ 원장님들의 가장 큰 착각", 3.2, 1, 8, 0.4, 18, bold=True, color=BLACK, align=PP_ALIGN.LEFT)

mistake_items = [
    '"예쁘게 만들면 된다"',
    '"고퀄 영상이면 된다"'
]
add_bullet_text(slide7, mistake_items, 3.7, 1.5, 7, 0.8, 16, color=GRAY)

# 경고 박스
warning_box = slide7.shapes.add_shape(1, Inches(1.5), Inches(4.7), Inches(7), Inches(0.8))
warning_box.fill.solid()
warning_box.fill.fore_color.rgb = RGBColor(220, 53, 69)
warning_box.line.fill.background()

add_text_box(slide7, "⚠️ 3초 안에 이해되지 않으면 끝입니다", 4.95, 2, 6, 0.4, 20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 핵심 원칙
principle_box = slide7.shapes.add_shape(1, Inches(1), Inches(5.8), Inches(8), Inches(1.2))
principle_box.fill.solid()
principle_box.fill.fore_color.rgb = BLACK
principle_box.line.fill.background()

add_text_box(slide7, "💡 크리에이티브는 랜딩으로 보내는 화살표입니다", 6.0, 1.5, 7, 0.4, 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide7, "여기서 모든 걸 설명하려 하면 광고도, 릴스도 망합니다", 6.4, 1.5, 7, 0.4, 16, color=WHITE, align=PP_ALIGN.CENTER)

# 슬라이드 8: 크리에이티브 (2) - 국룰 구조
slide8 = add_blank_slide()
slide8.background.fill.solid()
slide8.background.fill.fore_color.rgb = WHITE

title_bg = slide8.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = DEEP_NAVY
title_bg.line.fill.background()

add_title_box(slide8, "크리에이티브의 국룰 구조", 0.3, 1, 8, 0.6, 32, color=WHITE)

# 3단계 구조
steps = [
    ("1️⃣ 문제 직격", [
        '"광고비 쓰는데 문의 없는 분"',
        '"릴스 열심히 올리는데 예약 없는 분"'
    ]),
    ("2️⃣ 원인 반전", [
        '"광고가 아니라 구조 문제입니다"',
        '"릴스가 아니라 순서 문제입니다"'
    ]),
    ("3️⃣ 다음 행동 암시", [
        '"이걸 먼저 확인하세요"',
        '"이 구조가 없으면 광고는 의미 없습니다"'
    ])
]

y_position = 1.8
colors = [NAVY_LIGHT, RGBColor(59, 80, 137), RGBColor(79, 100, 157)]

for i, (step_title, items) in enumerate(steps):
    step_box = slide8.shapes.add_shape(1, Inches(1.2), Inches(y_position), Inches(7.6), Inches(1.4))
    step_box.fill.solid()
    step_box.fill.fore_color.rgb = colors[i]
    step_box.line.fill.background()

    add_text_box(slide8, step_title, y_position + 0.15, 1.5, 7, 0.4, 20, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    for j, item in enumerate(items):
        add_text_box(slide8, item, y_position + 0.55 + j * 0.35, 1.8, 6.8, 0.3, 15, color=WHITE, align=PP_ALIGN.LEFT)

    y_position += 1.6

# 슬라이드 9: 랜딩 (1)
slide9 = add_blank_slide()
slide9.background.fill.solid()
slide9.background.fill.fore_color.rgb = WHITE

title_bg = slide9.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = DEEP_NAVY
title_bg.line.fill.background()

add_title_box(slide9, "4️⃣ 랜딩 (Landing)", 0.3, 1, 8, 0.6, 36, color=WHITE)

# 정의 박스
def_box = slide9.shapes.add_shape(1, Inches(1), Inches(1.6), Inches(8), Inches(1.0))
def_box.fill.solid()
def_box.fill.fore_color.rgb = NAVY_LIGHT
def_box.line.fill.background()

add_text_box(slide9, "고객이 '결정하기 전에' 머무르는 공간", 1.9, 1.5, 7, 0.4, 20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide9, "❌ 원장님 현실", 3.0, 1, 8, 0.4, 18, bold=True, color=BLACK, align=PP_ALIGN.LEFT)

reality_items = [
    "릴스만 있음",
    "광고 눌렀는데 어디로 가야 할지 모름",
    "플레이스/프로필 정보 부족"
]
add_bullet_text(slide9, reality_items, 3.5, 1.5, 7, 1.2, 16, color=GRAY)

add_text_box(slide9, "→ 고객은 구경만 하고 나갑니다", 4.8, 1.5, 7, 0.4, 18, bold=True, color=BLACK, align=PP_ALIGN.LEFT)

# 진짜 역할 박스
role_box = slide9.shapes.add_shape(1, Inches(1), Inches(5.4), Inches(8), Inches(1.4))
role_box.fill.solid()
role_box.fill.fore_color.rgb = BLACK
role_box.line.fill.background()

add_text_box(slide9, "✅ 랜딩의 진짜 역할", 5.6, 1.5, 7, 0.4, 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide9, "설득 ❌  감동 ❌", 6.0, 1.5, 7, 0.4, 16, color=GRAY, align=PP_ALIGN.CENTER)
add_text_box(slide9, "정보 정리 + 불안 해소 ✅", 6.4, 1.5, 7, 0.4, 20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 슬라이드 10: 랜딩 (2) - 최소 구성
slide10 = add_blank_slide()
slide10.background.fill.solid()
slide10.background.fill.fore_color.rgb = WHITE

title_bg = slide10.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = DEEP_NAVY
title_bg.line.fill.background()

add_title_box(slide10, "랜딩의 최소 구성 (국룰)", 0.3, 1, 8, 0.6, 32, color=WHITE)

# 5가지 구성 요소
landing_elements = [
    "이 서비스가 누구에게 맞는지",
    "가격 범위 (정확하지 않아도 됨)",
    "진행 방식 (상담 → 시술)",
    "자주 묻는 질문 3~5개",
    "행동 버튼 (예약 / 상담)"
]

y_pos = 2.0
for i, element in enumerate(landing_elements):
    elem_box = slide10.shapes.add_shape(1, Inches(1.5), Inches(y_pos), Inches(7), Inches(0.65))
    elem_box.fill.solid()
    elem_box.fill.fore_color.rgb = NAVY_LIGHT if i % 2 == 0 else RGBColor(59, 80, 137)
    elem_box.line.fill.background()

    add_text_box(slide10, f"{i+1}. {element}", y_pos + 0.15, 1.8, 6.4, 0.4, 17, color=WHITE, align=PP_ALIGN.LEFT)
    y_pos += 0.8

# 하단 핵심 메시지
bottom_box = slide10.shapes.add_shape(1, Inches(1), Inches(6.2), Inches(8), Inches(0.9))
bottom_box.fill.solid()
bottom_box.fill.fore_color.rgb = BLACK
bottom_box.line.fill.background()

add_text_box(slide10, "💡 랜딩이 없는 광고는 사람을 데려와서", 6.3, 1.5, 7, 0.3, 16, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide10, "길바닥에 세워두는 것입니다", 6.6, 1.5, 7, 0.3, 16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 슬라이드 11: 전환 (1)
slide11 = add_blank_slide()
slide11.background.fill.solid()
slide11.background.fill.fore_color.rgb = WHITE

title_bg = slide11.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = DEEP_NAVY
title_bg.line.fill.background()

add_title_box(slide11, "5️⃣ 전환 (Conversion)", 0.3, 1, 8, 0.6, 36, color=WHITE)

# 정의 박스
def_box = slide11.shapes.add_shape(1, Inches(1), Inches(1.6), Inches(8), Inches(1.4))
def_box.fill.solid()
def_box.fill.fore_color.rgb = NAVY_LIGHT
def_box.line.fill.background()

add_text_box(slide11, "고객이 실제로 행동하게 만드는 장치", 1.9, 1.5, 7, 0.4, 20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide11, "고객의 귀찮음·불안·부담을", 2.4, 1.5, 7, 0.3, 16, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide11, "대신 처리해주는 시스템", 2.7, 1.5, 7, 0.3, 16, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide11, "❌ 전환이 안 되는 이유", 3.4, 1, 8, 0.4, 18, bold=True, color=BLACK, align=PP_ALIGN.LEFT)

no_conversion = [
    '"DM 주세요"',
    '"문의 주세요"'
]
add_bullet_text(slide11, no_conversion, 3.9, 1.5, 7, 0.8, 16, color=GRAY)

add_text_box(slide11, "→ 고객 입장에서는 너무 번거롭고, 너무 불안합니다", 4.8, 1.5, 7, 0.4, 17, bold=True, color=BLACK, align=PP_ALIGN.LEFT)

# 국룰 박스
rule_box = slide11.shapes.add_shape(1, Inches(1), Inches(5.5), Inches(8), Inches(1.4))
rule_box.fill.solid()
rule_box.fill.fore_color.rgb = BLACK
rule_box.line.fill.background()

add_text_box(slide11, "✅ 전환의 국룰", 5.7, 1.5, 7, 0.4, 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rules = "생각하지 않게 만들 것  |  망설이지 않게 만들 것  |  선택지를 단순화할 것"
add_text_box(slide11, rules, 6.2, 1.5, 7, 0.4, 15, color=WHITE, align=PP_ALIGN.CENTER)

# 슬라이드 12: 전환 (2) - 장치 예시
slide12 = add_blank_slide()
slide12.background.fill.solid()
slide12.background.fill.fore_color.rgb = WHITE

title_bg = slide12.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = DEEP_NAVY
title_bg.line.fill.background()

add_title_box(slide12, "전환 장치 예시", 0.3, 1, 8, 0.6, 32, color=WHITE)

# 4가지 장치
conversion_tools = [
    "상담 신청 폼",
    "예약 가능 날짜 표시",
    '"상담만 진행 가능" 문구',
    '"선착순 마감" 안내'
]

y_pos = 2.2
for i, tool in enumerate(conversion_tools):
    tool_box = slide12.shapes.add_shape(1, Inches(2), Inches(y_pos), Inches(6), Inches(0.75))
    tool_box.fill.solid()

    if i == 0 or i == 1:
        tool_box.fill.fore_color.rgb = NAVY_LIGHT
    else:
        tool_box.fill.fore_color.rgb = RGBColor(59, 80, 137)

    tool_box.line.fill.background()

    add_text_box(slide12, tool, y_pos + 0.2, 2.3, 5.4, 0.4, 19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    y_pos += 1.0

# 슬라이드 13: 구조의 본질
slide13 = add_blank_slide()
slide13.background.fill.solid()
slide13.background.fill.fore_color.rgb = DEEP_NAVY

add_title_box(slide13, "이 구조의 본질", 1.5, 1, 8, 0.8, 40, color=WHITE)

# 핵심 메시지 박스
essence_box = slide13.shapes.add_shape(1, Inches(1.5), Inches(2.8), Inches(7), Inches(2.2))
essence_box.fill.solid()
essence_box.fill.fore_color.rgb = WHITE
essence_box.line.fill.background()

add_text_box(slide13, "이 5단계 구조는", 3.0, 2, 6, 0.4, 22, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
add_text_box(slide13, "광고 전용 공식이 아니라", 3.5, 2, 6, 0.4, 22, color=BLACK, align=PP_ALIGN.CENTER)
add_text_box(slide13, "콘텐츠 제작 공식입니다", 4.0, 2, 6, 0.4, 22, bold=True, color=DEEP_NAVY, align=PP_ALIGN.CENTER)

add_text_box(slide13, "릴스, 블로그, 광고, 상담", 4.7, 2, 6, 0.4, 18, color=GRAY, align=PP_ALIGN.CENTER)

# 하단 강조 박스
bottom_emphasis = slide13.shapes.add_shape(1, Inches(1), Inches(5.8), Inches(8), Inches(1.2))
bottom_emphasis.fill.solid()
bottom_emphasis.fill.fore_color.rgb = BLACK
bottom_emphasis.line.fill.background()

add_text_box(slide13, "광고는 이 구조를", 6.0, 1.5, 7, 0.4, 20, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide13, "돈으로 증폭시키는 수단일 뿐입니다", 6.4, 1.5, 7, 0.4, 20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 슬라이드 14: 핵심 요약
slide14 = add_blank_slide()
slide14.background.fill.solid()
slide14.background.fill.fore_color.rgb = WHITE

title_bg = slide14.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = DEEP_NAVY
title_bg.line.fill.background()

add_title_box(slide14, "핵심 요약", 0.3, 1, 8, 0.6, 36, color=WHITE)

# 5단계 요약
summary_steps = [
    ("타겟", "실제로 돈을 쓰는 사람"),
    ("오퍼", "지금 해야 하는 이유"),
    ("크리에이티브", '"이건 내 얘기다" 3초 안에'),
    ("랜딩", "정보 정리 + 불안 해소"),
    ("전환", "생각·망설임 제거")
]

y_pos = 1.8
colors_cycle = [NAVY_LIGHT, RGBColor(59, 80, 137), RGBColor(79, 100, 157), RGBColor(39, 60, 117), DEEP_NAVY]

for i, (step, desc) in enumerate(summary_steps):
    step_box = slide14.shapes.add_shape(1, Inches(1.5), Inches(y_pos), Inches(7), Inches(0.75))
    step_box.fill.solid()
    step_box.fill.fore_color.rgb = colors_cycle[i]
    step_box.line.fill.background()

    add_text_box(slide14, f"{i+1}. {step}", y_pos + 0.1, 1.8, 2, 0.3, 18, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text_box(slide14, desc, y_pos + 0.35, 4, 4.2, 0.3, 16, color=WHITE, align=PP_ALIGN.LEFT)

    y_pos += 0.9

# 슬라이드 15: 마무리
slide15 = add_blank_slide()
slide15.background.fill.solid()
slide15.background.fill.fore_color.rgb = DEEP_NAVY

# 메인 메시지 박스
main_box = slide15.shapes.add_shape(1, Inches(1), Inches(2.2), Inches(8), Inches(2.8))
main_box.fill.solid()
main_box.fill.fore_color.rgb = WHITE
main_box.line.fill.background()

add_text_box(slide15, "광고가 어려운 게 아니라", 2.6, 1.5, 7, 0.5, 28, color=BLACK, align=PP_ALIGN.CENTER)
add_text_box(slide15, "광고 전에 필요한", 3.2, 1.5, 7, 0.5, 28, color=BLACK, align=PP_ALIGN.CENTER)
add_text_box(slide15, "구조가 없었던 겁니다", 3.8, 1.5, 7, 0.5, 32, bold=True, color=DEEP_NAVY, align=PP_ALIGN.CENTER)

# 저장
prs.save('/Users/jieun-/Desktop/바이브코딩/1월5일 무료라이브/광고_콘텐츠_국룰구조_PRD.pptx')

print("PPT 파일이 성공적으로 생성되었습니다!")
print("파일명: 광고_콘텐츠_국룰구조_PRD.pptx")
